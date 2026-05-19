import os
import shutil
import tempfile
import uuid
from pathlib import Path

from flask import Flask, jsonify, render_template, request, send_file, session

app = Flask(__name__)
app.secret_key = os.urandom(24)

UPLOAD_FOLDER = tempfile.mkdtemp(prefix="merger_")
ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".ppt"}
MAX_FILES = 20


def allowed(filename):
    return Path(filename).suffix.lower() in ALLOWED_EXTENSIONS


def session_dir(sid):
    d = Path(UPLOAD_FOLDER) / sid
    d.mkdir(parents=True, exist_ok=True)
    return d


@app.route("/")
def index():
    if "sid" not in session:
        session["sid"] = str(uuid.uuid4())
    return render_template("index1.html")


@app.route("/upload", methods=["POST"])
def upload():
    sid = session.get("sid", str(uuid.uuid4()))
    session["sid"] = sid
    sdir = session_dir(sid)

    order_file = sdir / "order.txt"
    files = order_file.read_text().splitlines() if order_file.exists() else []

    if len(files) >= MAX_FILES:
        return jsonify({"error": f"Maximum {MAX_FILES} files allowed."}), 400

    uploaded = request.files.getlist("files")
    added = []
    for f in uploaded:
        if len(files) >= MAX_FILES:
            break
        if not f.filename or not allowed(f.filename):
            continue
        safe_name = f"{uuid.uuid4().hex}_{Path(f.filename).name}"
        dest = sdir / safe_name
        f.save(dest)
        entry = f"{safe_name}|||{f.filename}"
        files.append(entry)
        added.append({"id": safe_name, "name": f.filename,
                      "ext": Path(f.filename).suffix.lower()})

    order_file.write_text("\n".join(files))
    return jsonify({"added": added, "total": len(files)})


@app.route("/files")
def list_files():
    sid = session.get("sid")
    if not sid:
        return jsonify({"files": []})
    sdir = session_dir(sid)
    order_file = sdir / "order.txt"
    if not order_file.exists():
        return jsonify({"files": []})
    result = []
    for line in order_file.read_text().splitlines():
        if "|||" not in line:
            continue
        safe, orig = line.split("|||", 1)
        result.append({"id": safe, "name": orig,
                        "ext": Path(orig).suffix.lower()})
    return jsonify({"files": result})


@app.route("/reorder", methods=["POST"])
def reorder():
    sid = session.get("sid")
    if not sid:
        return jsonify({"ok": False}), 400
    sdir = session_dir(sid)
    order_file = sdir / "order.txt"
    new_ids = request.json.get("ids", [])
    if not order_file.exists():
        return jsonify({"ok": False}), 400

    mapping = {}
    for line in order_file.read_text().splitlines():
        if "|||" in line:
            safe, orig = line.split("|||", 1)
            mapping[safe] = orig

    new_lines = [f"{nid}|||{mapping[nid]}" for nid in new_ids if nid in mapping]
    order_file.write_text("\n".join(new_lines))
    return jsonify({"ok": True})


@app.route("/remove", methods=["POST"])
def remove():
    sid = session.get("sid")
    if not sid:
        return jsonify({"ok": False}), 400
    sdir = session_dir(sid)
    order_file = sdir / "order.txt"
    file_id = request.json.get("id")
    if not order_file.exists():
        return jsonify({"ok": False}), 400

    lines = order_file.read_text().splitlines()
    new_lines = [l for l in lines if not l.startswith(file_id + "|||")]
    order_file.write_text("\n".join(new_lines))

    target = sdir / file_id
    if target.exists():
        target.unlink()
    return jsonify({"ok": True})


@app.route("/clear", methods=["POST"])
def clear():
    sid = session.get("sid")
    if not sid:
        return jsonify({"ok": True})
    sdir = session_dir(sid)
    if sdir.exists():
        shutil.rmtree(sdir)
    sdir.mkdir(parents=True, exist_ok=True)
    return jsonify({"ok": True})


@app.route("/merge", methods=["POST"])
def merge():
    sid = session.get("sid")
    if not sid:
        return jsonify({"error": "No session"}), 400
    sdir = session_dir(sid)
    order_file = sdir / "order.txt"
    if not order_file.exists():
        return jsonify({"error": "No files uploaded"}), 400

    entries = [(line.split("|||")[0], line.split("|||")[1])
               for line in order_file.read_text().splitlines() if "|||" in line]
    if len(entries) < 2:
        return jsonify({"error": "Upload at least 2 files"}), 400

    req_data = request.get_json(silent=True) or {}
    custom_name = req_data.get("filename", "").strip()
    if not custom_name:
        custom_name = "merged"

    paths = [str(sdir / safe) for safe, _ in entries]
    exts = {Path(orig).suffix.lower() for _, orig in entries}
    all_pdf  = exts <= {".pdf"}
    all_pptx = exts <= {".pptx", ".ppt"}

    try:
        if all_pdf:
            out, mime, dl_name = merge_pdfs(paths), "application/pdf", f"{custom_name}.pdf"
        elif all_pptx:
            out, mime, dl_name = (merge_pptx(paths),
                                  "application/vnd.openxmlformats-officedocument"
                                  ".presentationml.presentation",
                                  f"{custom_name}.pptx")
        else:
            pdf_paths = []
            tmpdir = tempfile.mkdtemp()
            try:
                for safe, orig in entries:
                    p = str(sdir / safe)
                    if Path(orig).suffix.lower() == ".pdf":
                        pdf_paths.append(p)
                    else:
                        pdf_paths.append(pptx_to_pdf(p, tmpdir))
                out = merge_pdfs(pdf_paths)
            finally:
                shutil.rmtree(tmpdir, ignore_errors=True)
            mime, dl_name = "application/pdf", f"{custom_name}.pdf"

        return send_file(out, mimetype=mime,
                         as_attachment=True, download_name=dl_name)
    except Exception as e:
        return jsonify({"error": str(e)}), 500


def merge_pdfs(paths):
    from pypdf import PdfWriter, PdfReader
    writer = PdfWriter()
    for p in paths:
        for page in PdfReader(p).pages:
            writer.add_page(page)
    out = tempfile.NamedTemporaryFile(suffix=".pdf", delete=False)
    writer.write(out)
    out.close()
    return out.name


def merge_pptx(paths):
    import copy
    from pptx import Presentation
    base = Presentation(paths[0])
    for p in paths[1:]:
        src = Presentation(p)
        for slide in src.slides:
            layout_idx = min(
                len(base.slide_layouts) - 1,
                src.slides.index(slide)
            )
            new_slide = base.slides.add_slide(base.slide_layouts[layout_idx])
            for ph in new_slide.placeholders:
                ph._element.getparent().remove(ph._element)
            for shape in slide.shapes:
                new_slide.shapes._spTree.insert(2, copy.deepcopy(shape._element))
    out = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    out.close()
    base.save(out.name)
    return out.name


def pptx_to_pdf(pptx_path, out_dir):
    import win32com.client
    import pythoncom
    
    pythoncom.CoInitialize()
    
    abs_pptx_path = os.path.abspath(pptx_path)
    output_pdf_path = os.path.abspath(os.path.join(out_dir, Path(pptx_path).stem + ".pdf"))
    
    powerpoint = None
    presentation = None
    try:
        powerpoint = win32com.client.DispatchEx("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(abs_pptx_path, WithWindow=False)
        presentation.SaveAs(output_pdf_path, FileFormat=32)
        
    except Exception as e:
        raise RuntimeError(f"Native Windows PowerPoint conversion failed: {str(e)}")
    finally:
        if presentation:
            presentation.Close()
        if powerpoint:
            powerpoint.Quit()
            
    if not os.path.exists(output_pdf_path):
        raise FileNotFoundError("PowerPoint completed execution, but the output PDF path is missing.")
        
    return output_pdf_path


if __name__ == "__main__":
    app.run(debug=True, port=5000)
