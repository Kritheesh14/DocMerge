from pptx import Presentation
import copy
from pathlib import Path

def merge_pptx(paths, output_path, status_callback=None):
    base = Presentation(paths[0])
    for i, p in enumerate(paths[1:], 1):
        if status_callback:
            status_callback(f"Merging {Path(p).name}...", (i / len(paths)) * 100)
            
        src = Presentation(p)
        for slide in src.slides:
            layout_idx = min(src.slides.index(slide), len(base.slide_layouts)-1)
            new_slide = base.slides.add_slide(base.slide_layouts[layout_idx])
            
            for ph in new_slide.placeholders:
                sp = ph._element
                sp.getparent().remove(sp)
                
            for shape in slide.shapes:
                el = copy.deepcopy(shape._element)
                new_slide.shapes._spTree.insert(2, el)
                
    base.save(output_path)
    return output_path
